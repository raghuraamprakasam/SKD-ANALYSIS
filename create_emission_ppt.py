import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.makedirs(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts",
    exist_ok=True,
)

df = pd.read_csv(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/COPY.csv"
)

pass_data = []
truck_data = []

# Passenger SKD (rows 1-18): Model = 'PASSENGER SKD' column, Emission = 'Unnamed: 6', Total = 'Unnamed: 7'
for i in range(1, 19):
    model = df.iloc[i]["PASSENGER SKD"]
    emission = df.iloc[i]["Unnamed: 6"]
    total = df.iloc[i]["Unnamed: 7"]
    if pd.notna(total) and not pd.isna(total):
        try:
            total_val = float(total)
            if total_val > 0:
                pass_data.append(
                    {
                        "Model": str(model) if pd.notna(model) else "Unknown",
                        "Emission": str(emission) if pd.notna(emission) else "N/A",
                        "Total": total_val,
                    }
                )
        except:
            pass

# Truck SKD (rows 20-43)
for i in range(20, 44):
    model = df.iloc[i]["PASSENGER SKD"]
    emission = df.iloc[i]["Unnamed: 6"]
    total = df.iloc[i]["Unnamed: 7"]
    if pd.notna(total) and not pd.isna(total):
        try:
            total_val = float(total)
            if total_val > 0:
                truck_data.append(
                    {
                        "Model": str(model) if pd.notna(model) else "Unknown",
                        "Emission": str(emission) if pd.notna(emission) else "N/A",
                        "Total": total_val,
                    }
                )
        except:
            pass

pass_df = pd.DataFrame(pass_data)
truck_df = pd.DataFrame(truck_data)

# Filter out total rows
pass_df = pass_df[pass_df["Model"] != "TOTAL PASS SKD"]
truck_df = truck_df[truck_df["Model"] != "TOTAL TRUCK SKD"]

print("=== PASSENGER SKD ===")
print(pass_df.to_string())
print(f"\nTotal Pass: {pass_df['Total'].sum():,.0f}")

print("\n=== TRUCK SKD ===")
print(truck_df.to_string())
print(f"\nTotal Truck: {truck_df['Total'].sum():,.0f}")

pass_total = pass_df["Total"].sum()
truck_total = truck_df["Total"].sum()

# === DOUGHNUT CHARTS ===
# 1. Passenger SKD by Emission
fig, ax = plt.subplots(figsize=(10, 8))
pass_emission = pass_df.groupby("Emission")["Total"].sum()
pass_emission = pass_emission.sort_values(ascending=False)
colors = ["#e74c3c", "#3498db", "#2ecc71", "#f39c12"]
wedges, texts, autotexts = ax.pie(
    pass_emission.values,
    labels=pass_emission.index,
    autopct="%1.1f%%",
    colors=colors[: len(pass_emission)],
    startangle=90,
    pctdistance=0.75,
    wedgeprops=dict(width=0.5, edgecolor="white"),
)
centre_circle = plt.Circle((0, 0), 0.35, fc="white")
ax.add_artist(centre_circle)
ax.set_title(
    f"PASSENGER SKD - By Emission Norm\n(Total: {pass_total:,.0f})",
    fontsize=14,
    fontweight="bold",
    pad=20,
)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/pass_emission_doughnut.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# 2. Truck SKD by Emission
fig, ax = plt.subplots(figsize=(10, 8))
truck_emission = truck_df.groupby("Emission")["Total"].sum()
truck_emission = truck_emission.sort_values(ascending=False)
colors = ["#9b59b6", "#1abc9c"]
wedges, texts, autotexts = ax.pie(
    truck_emission.values,
    labels=truck_emission.index,
    autopct="%1.1f%%",
    colors=colors[: len(truck_emission)],
    startangle=90,
    pctdistance=0.75,
    wedgeprops=dict(width=0.5, edgecolor="white"),
)
centre_circle = plt.Circle((0, 0), 0.35, fc="white")
ax.add_artist(centre_circle)
ax.set_title(
    f"TRUCK SKD - By Emission Norm\n(Total: {truck_total:,.0f})",
    fontsize=14,
    fontweight="bold",
    pad=20,
)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/truck_emission_doughnut.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# 3. HIERARCHICAL SUNBURST
fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(aspect="equal"))

sizes = [pass_total, truck_total]
categories = ["Passenger SKD", "Truck SKD"]
outer_colors = ["#3498db", "#e74c3c"]
outer_wedges, outer_texts, outer_autotexts = ax.pie(
    sizes,
    labels=categories,
    radius=1,
    colors=outer_colors,
    wedgeprops=dict(width=0.3, edgecolor="white"),
    textprops={"fontsize": 14, "fontweight": "bold"},
    autopct="%1.1f%%",
    pctdistance=0.85,
)

pass_emission_vals = pass_df.groupby("Emission")["Total"].sum()
truck_emission_vals = truck_df.groupby("Emission")["Total"].sum()

inner_radius = 0.7
colors_pass = ["#154360", "#1a5276", "#2874a6", "#3498db"]
colors_truck = ["#78281f", "#943126", "#b03a2e", "#e74c3c"]

for i, (label, val) in enumerate(pass_emission_vals.items()):
    if val > 0:
        ax.pie(
            [val],
            radius=inner_radius,
            colors=[colors_pass[i % len(colors_pass)]],
            wedgeprops=dict(width=0.15, edgecolor="white"),
        )

for i, (label, val) in enumerate(truck_emission_vals.items()):
    if val > 0:
        ax.pie(
            [val],
            radius=inner_radius,
            colors=[colors_truck[i % len(colors_truck)]],
            wedgeprops=dict(width=0.15, edgecolor="white"),
        )

ax.set_title(
    "HIERARCHICAL SUNBURST: SKD Type → Emission Norm",
    fontsize=14,
    fontweight="bold",
    pad=25,
)

plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/hierarchical_chart.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# === CREATE PPT ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.5))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(1.5))
tf = title_box.text_frame
tf.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "SKD PRODUCTION - EMISSION ANALYSIS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(1))
tf = subtitle.text_frame
tf.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = f"Passenger: {int(pass_total):,} | Truck: {int(truck_total):,}"
p.font.size = Pt(20)
p.font.color.rgb = WHITE

# Slide 2: Passenger Doughnut
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = RGBColor(52, 152, 219)

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[
    0
].text = "PASSENGER SKD - By Emission Norm (DOUGHNUT CHART)"
title.text_frame.paragraphs[0].font.size = Pt(20)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/pass_emission_doughnut.png",
    Inches(1),
    Inches(1.3),
    width=Inches(6),
)

text_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5), Inches(5))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "KEY INSIGHTS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(52, 152, 219)

insights = [f"Total: {int(pass_total):,} units"]
for emission, val in pass_emission_vals.items():
    pct = val / pass_total * 100
    insights.append(f"{emission}: {int(val):,} ({pct:.1f}%)")

for text in insights:
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(14)
    p.space_before = Pt(8)

# Slide 3: Truck Doughnut
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = RGBColor(231, 76, 60)

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "TRUCK SKD - By Emission Norm (DOUGHNUT CHART)"
title.text_frame.paragraphs[0].font.size = Pt(20)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/truck_emission_doughnut.png",
    Inches(1),
    Inches(1.3),
    width=Inches(6),
)

text_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5), Inches(5))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "KEY INSIGHTS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(231, 76, 60)

insights = [f"Total: {int(truck_total):,} units"]
for emission, val in truck_emission_vals.items():
    pct = val / truck_total * 100
    insights.append(f"{emission}: {int(val):,} ({pct:.1f}%)")

for text in insights:
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(14)
    p.space_before = Pt(8)

# Slide 4: Hierarchical
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "HIERARCHICAL VIEW - SKD Type & Emission Norm"
title.text_frame.paragraphs[0].font.size = Pt(20)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/new_charts/hierarchical_chart.png",
    Inches(0.5),
    Inches(1.2),
    width=Inches(12),
)

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(1))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = (
    "Outer Ring: SKD Type (Blue=Passenger, Red=Truck) | Inner Ring: Emission Standards"
)
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(100, 100, 100)
p.alignment = PP_ALIGN.CENTER

prs.save(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/SKD_Emission_Analysis.pptx"
)
print("\n[OK] PPT created: SKD_Emission_Analysis.pptx")
print(
    f"Total Pass: {int(pass_total):,} | Total Truck: {int(truck_total):,} | Grand: {int(pass_total + truck_total):,}"
)
