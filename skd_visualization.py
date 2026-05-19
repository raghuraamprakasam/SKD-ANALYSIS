import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

os.makedirs(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts",
    exist_ok=True,
)

xls = pd.ExcelFile(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/copy of rolling plan.xlsb"
)
df = pd.read_excel(xls, sheet_name="Sheet3")

pass_data = []
truck_data = []

for i in range(2, 19):
    row = df.iloc[i]
    model_num = row["Unnamed: 2"]
    stage = row["Unnamed: 3"]
    production = row["Unnamed: 5"]
    if pd.notna(model_num) and pd.notna(production):
        pass_data.append(
            {"Model_Number": model_num, "Stage": stage, "Total_Production": production}
        )

for i in range(20, 43):
    row = df.iloc[i]
    model_num = row["Unnamed: 2"]
    stage = row["Unnamed: 3"]
    production = row["Unnamed: 5"]
    if pd.notna(model_num) and pd.notna(production):
        truck_data.append(
            {"Model_Number": model_num, "Stage": stage, "Total_Production": production}
        )

pass_df = pd.DataFrame(pass_data)
truck_df = pd.DataFrame(truck_data)

# === PASSENGER SKD CHART ===
fig, ax = plt.subplots(figsize=(12, 8))
pass_sorted = pass_df.sort_values("Total_Production", ascending=True)
colors = plt.cm.Blues(np.linspace(0.3, 1, len(pass_sorted)))
bars = ax.barh(
    pass_sorted["Model_Number"].astype(str),
    pass_sorted["Total_Production"],
    color=colors,
)
ax.set_xlabel("Total Production (Units)", fontsize=12)
ax.set_ylabel("Model Number", fontsize=12)
ax.set_title(
    "PASSENGER SKD - Production by Model\n(Total: {:,} units)".format(
        int(pass_df["Total_Production"].sum())
    ),
    fontsize=14,
    fontweight="bold",
)
for i, v in enumerate(pass_sorted["Total_Production"]):
    ax.text(v + 30, i, str(int(v)), va="center", fontsize=9)
ax.set_xlim(0, pass_sorted["Total_Production"].max() * 1.15)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/pass_skd_by_model.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# === TRUCK SKD CHART ===
fig, ax = plt.subplots(figsize=(12, 8))
truck_sorted = truck_df.sort_values("Total_Production", ascending=True)
colors = plt.cm.Oranges(np.linspace(0.3, 1, len(truck_sorted)))
bars = ax.barh(
    truck_sorted["Model_Number"].astype(str),
    truck_sorted["Total_Production"],
    color=colors,
)
ax.set_xlabel("Total Production (Units)", fontsize=12)
ax.set_ylabel("Model Number", fontsize=12)
ax.set_title(
    "TRUCK SKD - Production by Model\n(Total: {:,} units)".format(
        int(truck_df["Total_Production"].sum())
    ),
    fontsize=14,
    fontweight="bold",
)
for i, v in enumerate(truck_sorted["Total_Production"]):
    ax.text(v + 10, i, str(int(v)), va="center", fontsize=9)
ax.set_xlim(0, truck_sorted["Total_Production"].max() * 1.15)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/truck_skd_by_model.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# === STAGE WISE PIE CHARTS ===
fig, axes = plt.subplots(1, 2, figsize=(14, 6))

pass_stage = pass_df.groupby("Stage")["Total_Production"].sum()
ax1 = axes[0]
colors1 = ["#2E86AB", "#A23B72", "#F18F01"]
ax1.pie(
    pass_stage.values,
    labels=pass_stage.index,
    autopct="%1.1f%%",
    colors=colors1[: len(pass_stage)],
    startangle=90,
    explode=[0.05] * len(pass_stage),
)
ax1.set_title(
    "PASS SKD - by Stage\n(Total: {:,})".format(int(pass_df["Total_Production"].sum())),
    fontsize=12,
    fontweight="bold",
)

truck_stage = truck_df.groupby("Stage")["Total_Production"].sum()
ax2 = axes[1]
ax2.pie(
    truck_stage.values,
    labels=truck_stage.index,
    autopct="%1.1f%%",
    colors=colors1[: len(truck_stage)],
    startangle=90,
    explode=[0.05] * len(truck_stage),
)
ax2.set_title(
    "TRUCK SKD - by Stage\n(Total: {:,})".format(
        int(truck_df["Total_Production"].sum())
    ),
    fontsize=12,
    fontweight="bold",
)

plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/skd_by_stage.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# === COMPARISON CHART ===
fig, ax = plt.subplots(figsize=(10, 6))
categories = ["Passenger SKD", "Truck SKD"]
totals = [pass_df["Total_Production"].sum(), truck_df["Total_Production"].sum()]
colors = ["#3498db", "#e74c3c"]
bars = ax.bar(categories, totals, color=colors, width=0.6)
ax.set_ylabel("Total Production (Units)", fontsize=12)
ax.set_title(
    "SKD Production Comparison\n(Total: {:,})".format(int(sum(totals))),
    fontsize=14,
    fontweight="bold",
)
for bar in bars:
    height = bar.get_height()
    ax.text(
        bar.get_x() + bar.get_width() / 2.0,
        height,
        f"{int(height):,}",
        ha="center",
        va="bottom",
        fontsize=14,
        fontweight="bold",
    )
ax.set_ylim(0, max(totals) * 1.15)
ax.axhline(y=sum(totals) / 2, color="gray", linestyle="--", alpha=0.5, label="Average")
plt.legend()
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/skd_comparison.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# === CREATE PPT ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "SKD PRODUCTION ANALYSIS"
subtitle.text = "Passenger SKD vs Truck SKD Comparison\nData Source: Rolling Plan"

# Slide 2: PASS SKD
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)
shapes = slide.shapes

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
tp = title_box.text_frame.paragraphs[0]
tp.text = "PASSENGER SKD - Production by Model"
tp.font.size = Pt(28)
tp.font.bold = True
tp.font.color.rgb = RGBColor(0, 51, 153)

img = shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/pass_skd_by_model.png",
    Inches(0.5),
    Inches(1.0),
    width=Inches(6),
)

text_box = shapes.add_textbox(Inches(7), Inches(1.0), Inches(5.5), Inches(5))
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "INTERPRETATION"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 153)

pass_total = int(pass_df["Total_Production"].sum())
pass_max = pass_df.loc[pass_df["Total_Production"].idxmax()]
pass_stage3 = pass_df[pass_df["Stage"] == "STAGE III"]["Total_Production"].sum()
pass_pct = pass_stage3 / pass_total * 100

summary = (
    f"\n• Total Production: {pass_total:,} units\n\n"
    f"• Highest: {pass_max['Model_Number']} ({int(pass_max['Total_Production'])} units)\n\n"
    f"• Stage Distribution:\n"
    f"  - Stage I: {int(pass_df[pass_df['Stage'] == 'STAGE I']['Total_Production'].sum())} ({pass_df[pass_df['Stage'] == 'STAGE I']['Total_Production'].sum() / pass_total * 100:.1f}%)\n"
    f"  - Stage II: {int(pass_df[pass_df['Stage'] == 'STAGE II']['Total_Production'].sum())} ({pass_df[pass_df['Stage'] == 'STAGE II']['Total_Production'].sum() / pass_total * 100:.1f}%)\n"
    f"  - Stage III: {int(pass_stage3)} ({pass_pct:.1f}%)\n\n"
    f"• Key Insight: {pass_pct:.1f}% of Passenger SKD is Stage III (Complete Build)"
)
p2 = tf.add_paragraph()
p2.text = summary
p2.font.size = Pt(12)

# Slide 3: TRUCK SKD
slide = prs.slides.add_slide(blank_layout)
shapes = slide.shapes

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
tp = title_box.text_frame.paragraphs[0]
tp.text = "TRUCK SKD - Production by Model"
tp.font.size = Pt(28)
tp.font.bold = True
tp.font.color.rgb = RGBColor(153, 51, 0)

img = shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/truck_skd_by_model.png",
    Inches(0.5),
    Inches(1.0),
    width=Inches(6),
)

text_box = shapes.add_textbox(Inches(7), Inches(1.0), Inches(5.5), Inches(5))
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "INTERPRETATION"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(153, 51, 0)

truck_total = int(truck_df["Total_Production"].sum())
truck_max = truck_df.loc[truck_df["Total_Production"].idxmax()]
truck_stage3 = truck_df[truck_df["Stage"] == "STAGE III"]["Total_Production"].sum()
truck_pct = truck_stage3 / truck_total * 100

summary = (
    f"\n• Total Production: {truck_total:,} units\n\n"
    f"• Highest: {truck_max['Model_Number']} ({int(truck_max['Total_Production'])} units)\n\n"
    f"• Stage Distribution:\n"
    f"  - Stage II: {int(truck_df[truck_df['Stage'] == 'STAGE II']['Total_Production'].sum())} ({truck_df[truck_df['Stage'] == 'STAGE II']['Total_Production'].sum() / truck_total * 100:.1f}%)\n"
    f"  - Stage III: {int(truck_stage3)} ({truck_pct:.1f}%)\n\n"
    f"• Key Insight: {truck_pct:.1f}% of Truck SKD is Stage III (Complete Build)\n\n"
    f"• No Stage I trucks in current plan"
)
p2 = tf.add_paragraph()
p2.text = summary
p2.font.size = Pt(12)

# Slide 4: COMBINED COMPARISON
slide = prs.slides.add_slide(blank_layout)
shapes = slide.shapes

title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
tp = title_box.text_frame.paragraphs[0]
tp.text = "COMBINED ANALYSIS - SKD PRODUCTION COMPARISON"
tp.font.size = Pt(22)
tp.font.bold = True
tp.font.color.rgb = RGBColor(51, 51, 51)

img1 = shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/skd_comparison.png",
    Inches(0.5),
    Inches(0.8),
    width=Inches(5.5),
)
img2 = shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/charts/skd_by_stage.png",
    Inches(6.2),
    Inches(0.8),
    width=Inches(6.5),
)

text_box = shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.8))
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "FINAL SUMMARY & INTERPRETATION"
p.font.size = Pt(16)
p.font.bold = True

total_all = pass_total + truck_total
summary = (
    f"\n• Total SKD Production: {total_all:,} units\n"
    f"• Passenger SKD: {pass_total:,} units ({pass_total / total_all * 100:.1f}%) | Truck SKD: {truck_total:,} units ({truck_total / total_all * 100:.1f}%)\n"
    f"• Passenger SKD is {pass_total / truck_total:.1f}x larger than Truck SKD\n"
    f"• Both categories dominated by Stage III (Pass: {pass_pct:.1f}%, Truck: {truck_pct:.1f}%)\n"
    f"• Passenger SKD has more model variants ({len(pass_df)} models) compared to Truck SKD ({len(truck_df)} models)"
)
p2 = tf.add_paragraph()
p2.text = summary
p2.font.size = Pt(12)

prs.save(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/python practise/SKD_Analysis_Presentation.pptx"
)
print("PPT created: SKD_Analysis_Presentation.pptx")
print(
    f"\nTotal Pass SKD: {pass_total:,} | Total Truck SKD: {truck_total:,} | Combined: {total_all:,}"
)
