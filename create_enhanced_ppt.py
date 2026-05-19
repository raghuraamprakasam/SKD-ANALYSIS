import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.makedirs(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts",
    exist_ok=True,
)

xls = pd.ExcelFile(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/copy of rolling plan.xlsb"
)
df = pd.read_excel(xls, sheet_name="Sheet3")

pass_data = []
truck_data = []

for i in range(2, 19):
    row = df.iloc[i]
    if pd.notna(row["Unnamed: 2"]) and pd.notna(row["Unnamed: 5"]):
        pass_data.append(
            {
                "Model_Number": str(row["Unnamed: 2"]),
                "Stage": str(row["Unnamed: 3"]),
                "Total_Production": row["Unnamed: 5"],
            }
        )

for i in range(20, 43):
    row = df.iloc[i]
    if pd.notna(row["Unnamed: 2"]) and pd.notna(row["Unnamed: 5"]):
        truck_data.append(
            {
                "Model_Number": str(row["Unnamed: 2"]),
                "Stage": str(row["Unnamed: 3"]),
                "Total_Production": row["Unnamed: 5"],
            }
        )

pass_df = pd.DataFrame(pass_data)
truck_df = pd.DataFrame(truck_data)

pass_total = int(pass_df["Total_Production"].sum())
truck_total = int(truck_df["Total_Production"].sum())
grand_total = pass_total + truck_total

# === CHARTS ===
# 1. Top 10 comparison
fig, ax = plt.subplots(figsize=(14, 8))
pass_sorted = pass_df.sort_values("Total_Production", ascending=False).head(10)
truck_sorted = truck_df.sort_values("Total_Production", ascending=False).head(10)
x = np.arange(10)
width = 0.35

colors_pass = plt.cm.Blues(np.linspace(0.5, 0.9, 10))
colors_truck = plt.cm.Oranges(np.linspace(0.5, 0.9, 10))

bars1 = ax.bar(
    x - width / 2,
    pass_sorted["Total_Production"],
    width,
    label=f"Passenger SKD ({pass_total:,})",
    color=colors_pass,
    edgecolor="white",
    linewidth=1.5,
)
bars2 = ax.bar(
    x + width / 2,
    truck_sorted["Total_Production"],
    width,
    label=f"Truck SKD ({truck_total:,})",
    color=colors_truck,
    edgecolor="white",
    linewidth=1.5,
)

ax.set_ylabel("Production (Units)", fontsize=12, fontweight="bold")
ax.set_title(
    "TOP 10 MODELS - PASSENGER vs TRUCK SKD COMPARISON",
    fontsize=14,
    fontweight="bold",
    pad=15,
)
ax.set_xticks(x)
ax.set_xticklabels([f"M{i + 1}" for i in range(10)], fontsize=10)
ax.legend(fontsize=11, loc="upper right", framealpha=0.9)
ax.grid(axis="y", alpha=0.3, linestyle="--")

for bar in bars1:
    ax.annotate(
        f"{int(bar.get_height())}",
        xy=(bar.get_x() + bar.get_width() / 2, bar.get_height()),
        xytext=(0, 3),
        textcoords="offset points",
        ha="center",
        va="bottom",
        fontsize=9,
        fontweight="bold",
    )
for bar in bars2:
    ax.annotate(
        f"{int(bar.get_height())}",
        xy=(bar.get_x() + bar.get_width() / 2, bar.get_height()),
        xytext=(0, 3),
        textcoords="offset points",
        ha="center",
        va="bottom",
        fontsize=9,
        fontweight="bold",
    )

plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/top10_comparison.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# 2. Stage distribution
fig, axes = plt.subplots(1, 2, figsize=(12, 6))
pass_stage = pass_df.groupby("Stage")["Total_Production"].sum()
truck_stage = truck_df.groupby("Stage")["Total_Production"].sum()

colors_stage = ["#e74c3c", "#3498db", "#2ecc71"]
axes[0].pie(
    pass_stage.values,
    labels=pass_stage.index,
    autopct="%1.1f%%",
    colors=colors_stage[: len(pass_stage)],
    startangle=90,
    explode=[0.05] * len(pass_stage),
)
axes[0].set_title(
    f"PASSENGER SKD STAGE\n(Total: {pass_total:,})", fontsize=12, fontweight="bold"
)
axes[1].pie(
    truck_stage.values,
    labels=truck_stage.index,
    autopct="%1.1f%%",
    colors=colors_stage[: len(truck_stage)],
    startangle=90,
    explode=[0.05] * len(truck_stage),
)
axes[1].set_title(
    f"TRUCK SKD STAGE\n(Total: {truck_total:,})", fontsize=12, fontweight="bold"
)

plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/stage_distribution.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# 3. Total distribution
fig, ax = plt.subplots(figsize=(8, 8))
sizes = [pass_total, truck_total]
labels = [
    f"PASSENGER SKD\n{pass_total:,}\n({pass_total / grand_total * 100:.1f}%)",
    f"TRUCK SKD\n{truck_total:,}\n({truck_total / grand_total * 100:.1f}%)",
]
colors = ["#3498db", "#e74c3c"]
explode = (0.05, 0.1)
wedges, texts = ax.pie(
    sizes,
    labels=labels,
    colors=colors,
    explode=explode,
    startangle=90,
    textprops={"fontsize": 12, "fontweight": "bold"},
)
ax.set_title(
    f"TOTAL SKD: {grand_total:,} UNITS", fontsize=14, fontweight="bold", pad=20
)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/total_distribution.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# 4. All models
fig, ax = plt.subplots(figsize=(12, 14))
combined = pd.concat([pass_df, truck_df])
combined = combined.sort_values("Total_Production", ascending=True)
colors_list = [
    "#3498db" if i < len(pass_df) else "#e74c3c" for i in range(len(combined))
]
bars = ax.barh(
    range(len(combined)),
    combined["Total_Production"],
    color=colors_list,
    edgecolor="white",
    linewidth=0.5,
)

ax.set_yticks(range(len(combined)))
ax.set_yticklabels(combined["Model_Number"], fontsize=8)
ax.set_xlabel("Total Production (Units)", fontsize=12, fontweight="bold")
ax.set_title("ALL SKD MODELS", fontsize=14, fontweight="bold", pad=15)
ax.grid(axis="x", alpha=0.3, linestyle="--")

for bar, val in zip(bars, combined["Total_Production"]):
    ax.text(
        val + 30,
        bar.get_y() + bar.get_height() / 2,
        f"{int(val):,}",
        va="center",
        fontsize=8,
        fontweight="bold",
    )

ax.set_xlim(0, combined["Total_Production"].max() * 1.18)
ax.legend(
    handles=[
        plt.Rectangle(
            (0, 0), 1, 1, facecolor="#3498db", label=f"Passenger ({pass_total:,})"
        ),
        plt.Rectangle(
            (0, 0), 1, 1, facecolor="#e74c3c", label=f"Truck ({truck_total:,})"
        ),
    ],
    loc="lower right",
)
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/all_models.png",
    dpi=150,
    bbox_inches="tight",
    facecolor="white",
)
plt.close()

# === CREATE PPT ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0, 85, 153)
DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.5))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
)
tf = title_box.text_frame
tf.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "SKD PRODUCTION ANALYSIS"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
tf = subtitle.text_frame
tf.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "Passenger SKD vs Truck SKD\nData Source: Rolling Plan"
p.font.size = Pt(20)
p.font.color.rgb = WHITE

footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
tf = footer.text_frame
tf.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = f"Total: {grand_total:,} Units | Pass: {pass_total:,} | Truck: {truck_total:,}"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(200, 200, 200)

# Slide 2: Executive Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "EXECUTIVE SUMMARY"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

# Cards
for idx, (label, value, pct, color) in enumerate(
    [
        ("TOTAL", grand_total, "100%", DARK_BLUE),
        (
            "PASSENGER",
            pass_total,
            f"{pass_total / grand_total * 100:.1f}%",
            RGBColor(52, 152, 219),
        ),
        (
            "TRUCK",
            truck_total,
            f"{truck_total / grand_total * 100:.1f}%",
            RGBColor(231, 76, 60),
        ),
    ]
):
    card = slide.shapes.add_shape(
        1, Inches(0.5 + idx * 4.3), Inches(1.3), Inches(4), Inches(2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = WHITE

    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = f"{value:,}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = pct
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

insights = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(3.5))
tf = insights.text_frame
p = tf.paragraphs[0]
p.text = "KEY INSIGHTS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

for text in [
    f"Passenger SKD dominates with {pass_total:,} units ({pass_total / grand_total * 100:.1f}%)",
    f"Truck SKD accounts for {truck_total:,} units ({truck_total / grand_total * 100:.1f}%)",
    f"Passenger to Truck ratio is {pass_total / truck_total:.1f}:1",
    "Stage III dominates in both categories (~94%)",
]:
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(13)
    p.space_before = Pt(8)

# Slide 3: Top 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "TOP 10 MODELS - COMPARISON"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/top10_comparison.png",
    Inches(0.5),
    Inches(1.2),
    width=Inches(12),
)

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.5))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = f"INTERPRETATION"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

p2 = tf.add_paragraph()
p2.text = f"• Highest Pass SKD: V6256809 (1,800) | Highest Truck SKD: V6553403 (360)\n• Top 10 Pass contribute {pass_sorted['Total_Production'].sum() / pass_total * 100:.1f}% | Top 10 Truck contribute {truck_sorted['Total_Production'].sum() / truck_total * 100:.1f}%"
p2.font.size = Pt(11)

# Slide 4: Stage Distribution
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "STAGE-WISE DISTRIBUTION"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/stage_distribution.png",
    Inches(0.5),
    Inches(1.2),
    width=Inches(6),
)
slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/total_distribution.png",
    Inches(7),
    Inches(1.2),
    width=Inches(5.5),
)

pass_s3 = pass_df[pass_df["Stage"] == "STAGE III"]["Total_Production"].sum()
truck_s3 = truck_df[truck_df["Stage"] == "STAGE III"]["Total_Production"].sum()

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.8))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "INTERPRETATION"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

p2 = tf.add_paragraph()
p2.text = f"• Stage III: Pass {pass_s3 / pass_total * 100:.1f}% | Truck {truck_s3 / truck_total * 100:.1f}%\n• Stage II significant in Truck SKD"
p2.font.size = Pt(11)

# Slide 5: All Models
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "ALL SKD MODELS - COMPLETE VIEW"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

slide.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/charts/all_models.png",
    Inches(0.5),
    Inches(1.1),
    width=Inches(12),
)

# Slide 6: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
title.text_frame.paragraphs[0].text = "CONCLUSION"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = WHITE

content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(5.5))
tf = content.text_frame
tf.word_wrap = True

conclusions = [
    (
        "Production Focus",
        f"Passenger SKD ({pass_total:,} units) - {pass_total / grand_total * 100:.1f}% of total",
    ),
    ("Stage Optimization", "Stage III dominates (~94%) - complete build preference"),
    ("Model Concentration", "Top 5 models contribute majority of production"),
    ("Truck Potential", "Truck SKD - 16.1% of total, growth potential"),
    ("Data Source", "Rolling Plan - Sheet3 (IO SKD models)"),
]

for title, desc in conclusions:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.space_before = Pt(4)

prs.save(
    "C:/Users/raghu/OneDrive/Desktop/visual studio codings/SKD analysis/SKD_Analysis_Presentation.pptx"
)
print("Enhanced PPT created!")
print(
    f"Slides: 6 | Pass: {pass_total:,} | Truck: {truck_total:,} | Total: {grand_total:,}"
)
